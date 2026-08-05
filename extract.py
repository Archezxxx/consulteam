import os
from pptx import Presentation
import fitz

base = os.path.dirname(os.path.abspath(__file__))
folder = os.path.join(base, "information")
out_dir = os.path.join(base, "information_extracted")
os.makedirs(out_dir, exist_ok=True)

for file in os.listdir(folder):
    filepath = os.path.join(folder, file)
    outpath = os.path.join(out_dir, file + ".txt")
    
    text = ""
    try:
        if file.endswith('.pdf'):
            with fitz.open(filepath) as doc:
                for page in doc:
                    text += page.get_text() + "\n"
        elif file.endswith('.pptx'):
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        with open(outpath, 'w', encoding='utf-8') as f:
            f.write(text)
    except Exception as e:
        print(f"Failed {file}: {e}")
